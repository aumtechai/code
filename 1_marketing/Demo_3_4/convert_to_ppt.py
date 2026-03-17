import asyncio
import os
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt

async def convert_full_demo():
    base_dir = r"c:\Projects\AA\at\1_marketing\Demo_3_4"
    intro_url = "file:///" + os.path.join(base_dir, "intro_slides_v2.html").replace("\\", "/")
    app_url = "https://aumtech.ai"
    
    out_dir = os.path.join(base_dir, "v2", "ppt_assets")
    os.makedirs(out_dir, exist_ok=True)
    
    # Define scenes with (Title, Subtitle, Notes)
    # Notes are taken from DEMO_SCRIPT.md
    scenes = [
        # Intro Slides (Index 0-5)
        ("The Challenge", "Fragmented Academic Data", """Today, universities face an enormous challenge. Thousands of students — each with a unique academic journey... aumtech.ai changes that."""),
        ("The Cost of Inefficiency", "Administrative Overload", """Healthcare got diagnostic AI. Legal got contract automation. But the university has been left behind... until now."""),
        ("Meet Aura", "Agentic Student Success", """Meet aumtech.ai — the world's first Agentic Student Success Platform. An intelligent academic operating system that thinks alongside every student."""),
        ("Architecture", "EdNex & Aura", """Project EdNex functions as a DataStream conduit, securely normalizing legacy data into a unified cloud hub for Aura's intelligence layer."""),
        ("Privacy & Security", "Enterprise Grade AI", """Every request passes through our Aura Privacy Gateway, where sensitive student PII is automatically stripped and tokenized."""),
        ("Login Page", "Your AI Navigator Starts Here", """The sign-in experience is clean and professional. Students log in with their university credentials."""),
        
        # App Screens
        ("Personalized Dashboard", "Live Success Metrics", """Welcome to the Student Dashboard. The first thing a student sees is a warm, personalized greeting from Aura. The On-Track Score reflects real-time progress."""),
        ("Get Aura (AI Chat)", "Conversational Guidance", """At the heart of aumtech.ai is Aura — a conversational academic agent powered by Google Gemini. It knows your courses, grades, and calendar."""),
        ("Academic Roadmap", "Courses & Progression", """Under Academics, students get a clear view of every course. The Degree Roadmap shows the entire journey visually. No more mysterious holds."""),
        ("Tutoring Center", "Intelligent Triage", """Syncs directly with Canvas/Blackboard. Our AI analyzes triage notes, generates TA briefs, and load-balances sessions in seconds."""),
        ("Wellness & Productivity", "The Whole Student", """Wellness Check-ins personalize recommendations. The Study Timer provides focus tools integrated into your success metrics."""),
        ("Holds Center", "Transparent Resolution", """Shows every active hold in plain language: what it is, why it exists, and how much is owed. One click to resolve."""),
        ("Financial Aid Nexus", "Scholarship Matching", """Gemini analyzes your GPA, major, and interests against available scholarships. Draft personal statements in seconds."""),
        ("Institutional Analytics", "Dean's Dashboard", """Admins see a live risk dashboard. Launch outreach campaigns in a click. See tutoring demand and TA impact in real-time."""),
    ]

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page(viewport={"width": 1920, "height": 1080})
        
        # --- PHASE 1: Intro Slides ---
        print("Capturing Intro Slides...")
        await page.goto(intro_url)
        await page.evaluate("if(typeof autoAdvance !== 'undefined') autoAdvance = false;")
        for i in range(6):
            await page.evaluate(f"jumpToSlide({i})")
            await asyncio.sleep(1.5)
            path = os.path.join(out_dir, f"slide_{i}.png")
            await page.screenshot(path=path)
            scenes[i] = scenes[i] + (path,)

        # --- PHASE 2: App Screens ---
        print("Logging in...")
        await page.goto(f"{app_url}/login")
        await page.fill("input[name='username']", "daniel.garrett12@txu.edu")
        await page.fill("input[name='password']", "password123")
        await page.click("button[type='submit']")
        await page.wait_for_selector(".sidebar, .nav-item, .hero-card")
        await asyncio.sleep(2)

        # Dashboard
        scenes[6] = scenes[6] + (os.path.join(out_dir, "app_0.png"),)
        await page.screenshot(path=scenes[6][3])

        # Chat
        await page.click("text=Get Aura")
        await page.wait_for_selector("input[placeholder*='Type']", timeout=10000)
        await asyncio.sleep(1)
        scenes[7] = scenes[7] + (os.path.join(out_dir, "app_1.png"),)
        await page.screenshot(path=scenes[7][3])

        # Roadmap
        await page.click("text=Degree Roadmap")
        await asyncio.sleep(2)
        scenes[8] = scenes[8] + (os.path.join(out_dir, "app_2.png"),)
        await page.screenshot(path=scenes[8][3])

        # Tutoring
        await page.click("text=Tutoring Center")
        await asyncio.sleep(2)
        scenes[9] = scenes[9] + (os.path.join(out_dir, "app_3.png"),)
        await page.screenshot(path=scenes[9][3])

        # Wellness
        await page.click("text=Wellness")
        await asyncio.sleep(1.5)
        scenes[10] = scenes[10] + (os.path.join(out_dir, "app_4.png"),)
        await page.screenshot(path=scenes[10][3])

        # Holds
        await page.click("text=Holds") # Likely Holds & Alerts
        await asyncio.sleep(1.5)
        scenes[11] = scenes[11] + (os.path.join(out_dir, "app_5.png"),)
        await page.screenshot(path=scenes[11][3])

        # Financial Aid
        await page.click("text=Financial Nexus")
        await asyncio.sleep(2)
        scenes[12] = scenes[12] + (os.path.join(out_dir, "app_6.png"),)
        await page.screenshot(path=scenes[12][3])

        # Admin Panel
        await page.goto(f"{app_url}/admin")
        await asyncio.sleep(2)
        scenes[13] = scenes[13] + (os.path.join(out_dir, "app_7.png"),)
        await page.screenshot(path=scenes[13][3])

        await browser.close()

    # --- PHASE 3: Create PPT ---
    print("Creating Presentation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for title, subtitle, notes, img_path in scenes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Notes
        n_slide = slide.notes_slide
        n_slide.notes_text_frame.text = f"{title.upper()}\n{subtitle}\n\n{notes}"

    final_ppt = os.path.join(base_dir, "v2", "Aumtech_Pitch_Deck_Full.pptx")
    prs.save(final_ppt)
    print(f"DONE: {final_ppt}")

if __name__ == "__main__":
    asyncio.run(convert_full_demo())
