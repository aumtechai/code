import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_executive_deck():
    prs = Presentation()
    
    # ── THEME COLORS (Modern Academic Blue & Gold) ───────────────────────
    PRIMARY_BLUE = RGBColor(30, 41, 59) # Slate 800 (0x1e, 0x29, 0x3b)
    ACCENT_INDIGO = RGBColor(79, 70, 229) # Indigo 600 (0x4f, 0x46, 0xe5)
    TEXT_GRAY = RGBColor(71, 85, 105) # Slate 600 (0x47, 0x55, 0x69)
    WHITE = RGBColor(255, 255, 255)
    
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.foreground_color.rgb = color

    def add_title_to_slide(slide, text, subtitle=None):
        title_shape = slide.shapes.title
        title_shape.text = text
        title_text_frame = title_shape.text_frame
        title_text_frame.paragraphs[0].font.size = Pt(36)
        title_text_frame.paragraphs[0].font.bold = True
        title_text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_text_frame = subtitle_shape.text_frame
            subtitle_text_frame.paragraphs[0].font.size = Pt(20)
            subtitle_text_frame.paragraphs[0].font.color.rgb = TEXT_GRAY

    # ── SLIDE 1: TITLE ──────────────────────────────────────────────────
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Aura Intelligence"
    subtitle.text = "Elevating Student Success & Institutional Excellence\nPersonalized Academic Navigation at Scale"
    
    # ── SLIDE 2: THE CHALLENGE ──────────────────────────────────────────
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "The Modern Student Success Crisis")
    tf = slide.placeholders[1].text_frame
    tf.text = "Institutions face unprecedented pressure to improve retention and graduation rates:"
    p = tf.add_paragraph()
    p.text = "• Retention Gaps: Identifying at-risk students before they disengage."
    p = tf.add_paragraph()
    p.text = "• Resource Overload: Faculty and advisors are stretched thin."
    p = tf.add_paragraph()
    p.text = "• Wellness & Mental Health: Holistic student support is no longer optional."
    p = tf.add_paragraph()
    p.text = "• Fragmented Systems: SIS, LMS, and Student Life data living in silos."

    # ── SLIDE 3: OUR MISSION ────────────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "The Aura Mission")
    tf = slide.placeholders[1].text_frame
    tf.text = "Personalizing the Student Journey through Proactive AI Intelligence"
    p = tf.add_paragraph()
    p.text = "Aura is not just a chatbot; it's a unified academic ecosystem that connects institutional data with the student's daily experience."

    # ── SLIDE 4: EXECUTIVE SUMMARY ──────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "A Unified Solution for Campus Leadership")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• Proactive: Detects drops in engagement and GPA automatically."
    p = tf.add_paragraph()
    p.text = "• Integrative: Seamlessly verified middleware (EdNex) for SIS/LMS."
    p = tf.add_paragraph()
    p.text = "• Scalable: 24/7 support for 10,000+ students without increasing headcount."

    # ── SLIDE 5: PILLAR 1 - PROACTIVE GUIDANCE ──────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "Pillar 1: Proactive Academic Guidance")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Aura monitors the 'On-Track' score of every student."
    p = tf.add_paragraph()
    p.text = "• Early Warning System: Alerts advisors to critical GPA drops."
    p = tf.add_paragraph()
    p.text = "• SMART Course Load: Recommends schedules based on historic success paths."

    # ── SLIDE 6: PILLAR 2 - 24/7 AI NAVIGATOR ───────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "Pillar 2: The 24/7 AI Navigator")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Immediate, accurate answers for the student's entire academic life."
    p = tf.add_paragraph()
    p.text = "• Dynamic holds resolution."
    p = tf.add_paragraph()
    p.text = "• Financial aid status and scholarship matching."
    p = tf.add_paragraph()
    p.text = "• Course prerequisites and degree roadmap visualization."

    # ── SLIDE 7: PILLAR 3 - SPECIALIZED LEARNING ────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "Pillar 3: Specialized Learning Support")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Closing the loop between identifying a struggle and mastering the material."
    p = tf.add_paragraph()
    p.text = "• Instant Tutor Matching: Connecting students with peer and AI tutors."
    p = tf.add_paragraph()
    p.text = "• Syllabus Scanner: Breaking down complex course requirements."

    # ── SLIDE 8: PILLAR 4 - HOLISTIC WELLNESS ──────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "Pillar 4: Holistic Student Wellness")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "University Deans recognize that student success depends on wellness."
    p = tf.add_paragraph()
    p.text = "• Non-academic intervention (Financial, Social, Emotional)."
    p = tf.add_paragraph()
    p.text = "• Connection to Campus Resources: Directly booking wellness checks."

    # ── SLIDE 9: INSTITUTIONAL INSIGHTS ─────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "Institutional Intelligence for Leadership")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Data-driven decision making for the Dean's office."
    p = tf.add_paragraph()
    p.text = "• Aggregated Student Cohort Analytics."
    p = tf.add_paragraph()
    p.text = "• Departmental Performance heatmaps."
    p = tf.add_paragraph()
    p.text = "• Real-time demand forecasting for campus support services."

    # ── SLIDE 10: THE ADMINISTRATOR EXPERIENCE ──────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "The Administrator Control Center")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Empowering faculty to focus on teaching, not administrative routing."
    p = tf.add_paragraph()
    p.text = "• Bulk intervention triggers."
    p = tf.add_paragraph()
    p.text = "• Automated follow-up on critical student alerts."

    # ── SLIDE 11: SECURITY & INTEGRATION ────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "Architected for Institutional Security")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• EdNex Verified: Enterprise-grade SIS and LMS integration."
    p = tf.add_paragraph()
    p.text = "• FERPA compliant data handling."
    p = tf.add_paragraph()
    p.text = "• Cloud-Native Scalability with zero on-prem footprint."

    # ── SLIDE 12: RESULTS & ROI ─────────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "Measurable Institutional ROI")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• 15% increase in retention for targeted high-risk cohorts."
    p = tf.add_paragraph()
    p.text = "• 4x reduction in administrative 'bounce-around' for student queries."
    p = tf.add_paragraph()
    p.text = "• Faster time-to-graduation through optimized course pathway navigation."

    # ── SLIDE 13: VIDEO SHOWCASE (PLACEHOLDER) ──────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "9.5 Minute Master Platform Demo")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "[Insert Demo Video Here]"
    p = tf.add_paragraph()
    p.text = "A full walkthrough of the student experience on the live Aumtech platform."

    # ── SLIDE 14: IMPLEMENTATION ROADMAP ────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "The Roadmap to Adoption")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "1. Pilot Phase: Selected Department Focus (4 Weeks)."
    p = tf.add_paragraph()
    p.text = "2. SIS Integration: Secure data bridging."
    p = tf.add_paragraph()
    p.text = "3. Campus-Wide Launch & Adoption Monitoring."

    # ── SLIDE 15: CONCLUSION ────────────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    add_title_to_slide(slide, "Transforming Student Success Together")
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Let's innovate the future of the student experience."
    p = tf.add_paragraph()
    p.text = "Next Steps: Strategic Alignment & Pilot Scoping."

    output_path = r"C:\Projects\AA\at\1_marketing\Aura_Dean_Executive_Proposal.pptx"
    prs.save(output_path)
    print(f"Executive Deck created: {output_path}")

if __name__ == "__main__":
    create_executive_deck()
